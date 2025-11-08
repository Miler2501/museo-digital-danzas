#!/usr/bin/env python3
# create_presentation.py
# Genera Museo_Digital_Danzas_Peruanas.pptx con:
# - Menú interactivo (botones por danza)
# - Para cada danza: Historia, Vestimenta, Video (enlace a YouTube)
# - Botón "Volver al menú" en cada diapositiva de contenido
#
# Requisitos:
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Nombre de salida
OUTPUT = "Museo_Digital_Danzas_Peruanas.pptx"

# Lista de danzas (7)
DANZAS = [
    "Huayno",
    "Marinera",
    "Festejo",
    "Diablada",
    "Tondero",
    "Saya",
    "Carnaval"
]

# Ejemplo de URLs de video (reemplaza por las que prefieras)
VIDEO_LINKS = {
    "Huayno": "https://www.youtube.com/watch?v=PLACEHOLDER_HUAYNO",
    "Marinera": "https://www.youtube.com/watch?v=PLACEHOLDER_MARINERA",
    "Festejo": "https://www.youtube.com/watch?v=PLACEHOLDER_FESTEJO",
    "Diablada": "https://www.youtube.com/watch?v=PLACEHOLDER_DIABLADA",
    "Tondero": "https://www.youtube.com/watch?v=PLACEHOLDER_TONDERO",
    "Saya": "https://www.youtube.com/watch?v=PLACEHOLDER_SAYA",
    "Carnaval": "https://www.youtube.com/watch?v=PLACEHOLDER_CARNAVAL",
}

# Textos de ejemplo (puedes editarlos después en PowerPoint)
HISTORY_TEXT = {
    "Huayno": "El huayno es una expresión musical y dancística andina con raíces precolombinas y mestizas...",
    "Marinera": "La marinera es un baile costeño considerado un símbolo de elegancia y coqueteo...",
    "Festejo": "El festejo es una danza afroperuana de la costa peruana que celebra...",
    "Diablada": "La diablada es una danza de fuerte carácter ritual originaria del Altiplano...",
    "Tondero": "El tondero es un baile costeño del norte del Perú con influencias españolas y africanas...",
    "Saya": "La saya es una expresión afroanda de la región del sur, con ritmos marcados...",
    "Carnaval": "Las danzas de carnaval varían por región; suelen incluir comparsas, máscaras y música festiva..."
}

CLOTHING_TEXT = {
    "Huayno": "- Mujer: pollera, manta bordada, chompa de lana, sombrero.\n- Hombre: poncho, sombrero, pantalón resistente.",
    "Marinera": "- Mujer: vestido elegante, pollera, pañuelo blanco.\n- Hombre: traje claro, sombrero, pañuelo.",
    "Festejo": "- Mujer: vestido con vuelo, accesorios de estación.\n- Hombre: camisa y pantalón con colores vivos.",
    "Diablada": "- Trajes con máscaras, colorido y detalles simbólicos.",
    "Tondero": "- Mujer: falda amplia, blusas bordadas.\n- Hombre: camisa, pantalón y sombrero típico.",
    "Saya": "- Ropa con influencia afro, adornos de cintura.\n- Hombre: vestimenta sencilla, elementos rituales.",
    "Carnaval": "- Varía por región: trajes coloridos, máscaras y adornos festivos."
}

# Funciones de ayuda
def add_title(slide, title):
    title_box = slide.shapes.title
    title_box.text = title
    title_tf = title_box.text_frame
    for paragraph in title_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True


def add_text_box(slide, left, top, width, height, text, font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = 'Calibri'


def make_button(slide, left, top, width, height, label, font_size=14, fill_rgb=(0,112,192)):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # 1 = Rectangle (MSO_SHAPE_RECTANGLE)
    )
    # Relleno color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*fill_rgb)
    # Texto
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255,255,255)
    return shape


def main():
    prs = Presentation()
    # Usa layout 0 para título, 1 para contenido
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # 1) Portada
    s = prs.slides.add_slide(title_layout)
    add_title(s, "Museo Digital de las Danzas Peruanas")
    add_text_box(s, Inches(1), Inches(2.2), Inches(8), Inches(1), "Tema: Identidad\nPresentación interactiva (E-Book).")

    # 2) Menú principal
    menu = prs.slides.add_slide(content_layout)
    add_title(menu, "Menú - Selecciona una danza")
    # Crear botones para cada danza y almacenar posición (en grid)
    buttons = {}
    left_start = Inches(0.5)
    top_start = Inches(2)
    btn_w = Inches(3)
    btn_h = Inches(0.7)
    gap_x = Inches(0.4)
    gap_y = Inches(0.3)
    cols = 2
    for idx, danza in enumerate(DANZAS):
        row = idx // cols
        col = idx % cols
        left = left_start + col*(btn_w + gap_x)
        top = top_start + row*(btn_h + gap_y)
        btn = make_button(menu, left, top, btn_w, btn_h, danza)
        buttons[danza] = btn

    # Guarda referencias de la diapositiva de menu
    menu_slide = menu

    # 3) Crear slides por danza: Historia, Vestimenta, Video
    # Mantener dict para la primera diapositiva de cada danza (para link desde menú)
    first_slide_of = {}

    for danza in DANZAS:
        # Historia
        s_hist = prs.slides.add_slide(content_layout)
        add_title(s_hist, f"{danza} — Historia")
        add_text_box(s_hist, Inches(0.5), Inches(1.9), Inches(9), Inches(3), HISTORY_TEXT.get(danza, ""))
        # Botón Volver al menú
        btn_return = make_button(s_hist, Inches(8), Inches(0.2), Inches(1.2), Inches(0.5), "Inicio", fill_rgb=(80,80,80))
        # Vestimenta
        s_vest = prs.slides.add_slide(content_layout)
        add_title(s_vest, f"{danza} — Vestimenta")
        add_text_box(s_vest, Inches(0.5), Inches(1.9), Inches(9), Inches(3), CLOTHING_TEXT.get(danza, ""))
        btn_return2 = make_button(s_vest, Inches(8), Inches(0.2), Inches(1.2), Inches(0.5), "Inicio", fill_rgb=(80,80,80))
        # Video (aquí ponemos un enlace externo "Ver video")
        s_vid = prs.slides.add_slide(content_layout)
        add_title(s_vid, f"{danza} — Video")
        add_text_box(s_vid, Inches(0.5), Inches(1.9), Inches(9), Inches(1.2), "Haz clic en 'Ver video' para abrir el enlace en tu navegador (requiere conexión a Internet).")
        btn_video = make_button(s_vid, Inches(2.5), Inches(3), Inches(4), Inches(0.8), "Ver video", fill_rgb=(220,20,60))
        btn_return3 = make_button(s_vid, Inches(8), Inches(0.2), Inches(1.2), Inches(0.5), "Inicio", fill_rgb=(80,80,80))

        # Guardar referencia al primer slide (Historia) para enlazar desde menú
        first_slide_of[danza] = s_hist

        # Configurar enlaces: Volver al menú (target_slide = menu_slide)
        # python-pptx soporta click_action.target_slide
        btn_return.click_action.target_slide = menu_slide
        btn_return2.click_action.target_slide = menu_slide
        btn_return3.click_action.target_slide = menu_slide

        # Enlace del botón de video a YouTube (enlace externo)
        video_url = VIDEO_LINKS.get(danza, "")
        if video_url:
            btn_video.click_action.hyperlink.address = video_url

    # Enlazar los botones del menú a la primera diapositiva de cada danza
    for danza, btn in buttons.items():
        target = first_slide_of.get(danza)
        if target:
            btn.click_action.target_slide = target

    # Guardar archivo
    prs.save(OUTPUT)
    print(f"Presentación generada: {OUTPUT}")

if __name__ == "__main__":
    main()